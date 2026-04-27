#!/usr/bin/env python3
"""
Generador de documentos Word y PowerPoint para Project Coordinator Hub.

Uso:
  python3 tools/generar_documentos.py --proyecto gama-leasing --tipo minuta
  python3 tools/generar_documentos.py --proyecto gama-leasing --tipo resumen
  python3 tools/generar_documentos.py --proyecto gama-leasing --tipo presentacion

Tipos disponibles:
  minuta        -> Word (.docx) con plantilla de minuta
  resumen       -> Word (.docx) con resumen ejecutivo
  presentacion  -> PowerPoint (.pptx) con slides ejecutivos
"""

import argparse
import os
import csv
from datetime import date
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPtPt
from pptx.dml.color import RGBColor as PPTRGBColor

PROJECTS_DIR = Path(__file__).parent.parent / "projects"
OUTPUTS_DIR = Path(__file__).parent.parent / "outputs"

BRAND_DARK = RGBColor(0x1F, 0x48, 0x7E)     # azul oscuro
BRAND_ACCENT = RGBColor(0x00, 0xA9, 0xE0)   # azul claro
PPT_DARK = PPTRGBColor(0x1F, 0x48, 0x7E)
PPT_ACCENT = PPTRGBColor(0x00, 0xA9, 0xE0)


def find_project(nombre_parcial: str) -> Path:
    matches = [p for p in PROJECTS_DIR.iterdir() if p.is_dir() and nombre_parcial.lower() in p.name.lower()]
    if not matches:
        raise FileNotFoundError(f"No se encontro carpeta de proyecto con '{nombre_parcial}' en {PROJECTS_DIR}")
    if len(matches) > 1:
        print(f"Multiples coincidencias: {[m.name for m in matches]}")
        print(f"Usando: {matches[0].name}")
    return matches[0]


def leer_seguimiento(project_path: Path) -> list:
    f = project_path / "tareas" / "seguimiento.md"
    if not f.exists():
        return []
    lines = f.read_text(encoding="utf-8").splitlines()
    tareas = [l.strip() for l in lines if l.strip().startswith("- [")]
    return tareas


def leer_horas(project_path: Path) -> list:
    f = project_path / "horas" / "registro-horas.csv"
    if not f.exists():
        return []
    rows = []
    with open(f, encoding="utf-8") as fh:
        reader = csv.DictReader(fh)
        for row in reader:
            if row.get("persona"):
                rows.append(row)
    return rows


# ── WORD: MINUTA ─────────────────────────────────────────────────────────────

def generar_minuta(project_path: Path, output_path: Path):
    nombre_proyecto = project_path.name
    doc = Document()

    # Titulo
    titulo = doc.add_heading("Minuta de Reunion", level=0)
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = titulo.runs[0]
    run.font.color.rgb = BRAND_DARK

    # Metadatos
    doc.add_heading("Informacion del proyecto", level=1).runs[0].font.color.rgb = BRAND_DARK
    tabla_meta = doc.add_table(rows=5, cols=2)
    tabla_meta.style = "Table Grid"
    meta = [
        ("Proyecto", nombre_proyecto),
        ("Cliente", ""),
        ("Fecha", str(date.today())),
        ("Facilitador", ""),
        ("Participantes", ""),
    ]
    for i, (k, v) in enumerate(meta):
        tabla_meta.cell(i, 0).text = k
        tabla_meta.cell(i, 1).text = v

    # Secciones
    secciones = ["Objetivo", "Temas tratados", "Acuerdos (Accion | Responsable | Fecha)", "Riesgos y bloqueos", "Proximos pasos"]
    for s in secciones:
        doc.add_heading(s, level=1).runs[0].font.color.rgb = BRAND_DARK
        doc.add_paragraph("")

    # Firma
    doc.add_paragraph(f"\nDocumento generado el {date.today()} por Project Coordinator Hub")

    out = output_path / f"{date.today()}-minuta.docx"
    doc.save(out)
    print(f"Minuta guardada: {out}")


# ── WORD: RESUMEN EJECUTIVO ──────────────────────────────────────────────────

def generar_resumen(project_path: Path, output_path: Path):
    nombre_proyecto = project_path.name
    tareas = leer_seguimiento(project_path)
    horas = leer_horas(project_path)

    doc = Document()

    titulo = doc.add_heading("Resumen Ejecutivo", level=0)
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    titulo.runs[0].font.color.rgb = BRAND_DARK

    doc.add_heading("Datos del proyecto", level=1).runs[0].font.color.rgb = BRAND_DARK
    tabla = doc.add_table(rows=3, cols=2)
    tabla.style = "Table Grid"
    for i, (k, v) in enumerate([("Proyecto", nombre_proyecto), ("Fecha", str(date.today())), ("Estado general", "Verde / Amarillo / Rojo")]):
        tabla.cell(i, 0).text = k
        tabla.cell(i, 1).text = v

    for seccion in ["Avances clave", "Pendientes criticos", "Riesgos", "Decisiones requeridas", "Proxima semana"]:
        doc.add_heading(seccion, level=1).runs[0].font.color.rgb = BRAND_DARK
        doc.add_paragraph("")

    if tareas:
        doc.add_heading("Tareas en seguimiento", level=1).runs[0].font.color.rgb = BRAND_DARK
        for t in tareas:
            doc.add_paragraph(t, style="List Bullet")

    if horas:
        doc.add_heading("Registro de horas", level=1).runs[0].font.color.rgb = BRAND_DARK
        th = doc.add_table(rows=1 + len(horas), cols=4)
        th.style = "Table Grid"
        for j, h in enumerate(["Persona", "Rol", "Actividad", "Horas"]):
            th.cell(0, j).text = h
        for i, row in enumerate(horas):
            th.cell(i + 1, 0).text = row.get("persona", "")
            th.cell(i + 1, 1).text = row.get("rol", "")
            th.cell(i + 1, 2).text = row.get("actividad", "")
            th.cell(i + 1, 3).text = row.get("horas", "")

    doc.add_paragraph(f"\nDocumento generado el {date.today()} por Project Coordinator Hub")

    out = output_path / f"{date.today()}-resumen-ejecutivo.docx"
    doc.save(out)
    print(f"Resumen guardado: {out}")


# ── PPTX: PRESENTACION EJECUTIVA ─────────────────────────────────────────────

def generar_presentacion(project_path: Path, output_path: Path):
    nombre_proyecto = project_path.name
    tareas = leer_seguimiento(project_path)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(layout_idx=6):
        layout = prs.slide_layouts[layout_idx]
        return prs.slides.add_slide(layout)

    def titulo_slide(slide, titulo, subtitulo=""):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo
        run.font.size = PPtPt(32)
        run.font.bold = True
        run.font.color.rgb = PPT_DARK
        if subtitulo:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitulo
            r2.font.size = PPtPt(18)
            r2.font.color.rgb = PPT_ACCENT

    def cuerpo_slide(slide, items: list, top=Inches(1.6)):
        txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(5.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {item}"
            run.font.size = PPtPt(18)
            run.font.color.rgb = PPT_DARK

    # Slide 1: portada
    s1 = add_slide()
    titulo_slide(s1, nombre_proyecto, f"Presentacion ejecutiva  ·  {date.today()}")

    # Slide 2: de que trata
    s2 = add_slide()
    titulo_slide(s2, "Sobre el proyecto")
    cuerpo_slide(s2, [
        "Cliente: " + nombre_proyecto.split("-")[0].replace("-", " ").title(),
        "Codigo: " + nombre_proyecto,
        "Iniciativa: ver README del proyecto",
        "Fecha inicio: " + (nombre_proyecto[-10:] if len(nombre_proyecto) >= 10 else ""),
    ])

    # Slide 3: estado
    s3 = add_slide()
    titulo_slide(s3, "Estado del proyecto")
    cuerpo_slide(s3, ["Estado general: Por definir", "Avances clave: Por completar", "Riesgos activos: Por completar"])

    # Slide 4: tareas
    s4 = add_slide()
    titulo_slide(s4, "Seguimiento de tareas")
    cuerpo_slide(s4, tareas if tareas else ["Sin tareas registradas aun"])

    # Slide 5: proximos pasos
    s5 = add_slide()
    titulo_slide(s5, "Proximos pasos")
    cuerpo_slide(s5, ["Completar minuta de kickoff", "Registrar horas del equipo", "Definir riesgos y acuerdos"])

    # Slide 6: cierre
    s6 = add_slide()
    titulo_slide(s6, "Contacto y cierre", "Project Coordinator Hub · " + str(date.today()))

    out = output_path / f"{date.today()}-presentacion.pptx"
    prs.save(out)
    print(f"Presentacion guardada: {out}")


# ── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Genera documentos Word y PowerPoint por proyecto.")
    parser.add_argument("--proyecto", required=True, help="Nombre parcial de la carpeta del proyecto")
    parser.add_argument("--tipo", required=True, choices=["minuta", "resumen", "presentacion"], help="Tipo de documento a generar")
    args = parser.parse_args()

    project_path = find_project(args.proyecto)
    output_path = project_path / ("minutas" if args.tipo == "minuta" else "resumenes" if args.tipo == "resumen" else "resumenes")
    OUTPUTS_DIR.mkdir(exist_ok=True)
    output_path.mkdir(parents=True, exist_ok=True)

    if args.tipo == "minuta":
        generar_minuta(project_path, output_path)
    elif args.tipo == "resumen":
        generar_resumen(project_path, output_path)
    elif args.tipo == "presentacion":
        generar_presentacion(project_path, output_path)


if __name__ == "__main__":
    main()
